# -*- coding: utf-8 -*-
"""Genere 3 slides PPTX 'Anacoluthe' a inserer dans le brief Marseillan-Sete.
Format 12192000 x 6858000 EMU (13,333 x 7,5 pouces) = identique au deck de depart.
Style Anacoluthe : fond creme, header rouge UPPERCASE, sections teal, zero ombre.
Lance depuis la racine du repo :  python promo/build_slides_marseillan.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---------- Palette ----------
BRICK    = RGBColor(0xC6, 0x28, 0x28)
TEAL     = RGBColor(0x3D, 0x8B, 0x87)
AMBER    = RGBColor(0xC9, 0x6A, 0x30)
NAVY     = RGBColor(0x1E, 0x3A, 0x5F)
BLEU     = RGBColor(0x00, 0x66, 0xAD)
CREME    = RGBColor(0xFF, 0xFD, 0xF9)
BLOCBG   = RGBColor(0xF8, 0xFA, 0xFC)
TEAL50   = RGBColor(0xE8, 0xF4, 0xF3)
AMBER50  = RGBColor(0xFF, 0xF8, 0xF0)
NAVY100  = RGBColor(0xC9, 0xD5, 0xE3)
TEXTDARK = RGBColor(0x2D, 0x37, 0x48)
TEXTMUT  = RGBColor(0x71, 0x80, 0x96)
TEXTLT   = RGBColor(0xA0, 0xA8, 0xB5)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LINELT   = RGBColor(0xEC, 0xED, 0xF0)

HEAD = "Merriweather Sans"
BODY = "Merriweather"

IMG = "assets/images/"

prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]

SW, SH = 13.3333, 7.5   # pouces

# ---------- Helpers ----------
def no_shadow(shape):
    shape.shadow.inherit = False

def bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = CREME
    r.line.fill.background(); no_shadow(r)
    return r

def rrect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    no_shadow(s)
    return s

def hline(slide, x1, y1, x2, y2, color=LINELT, w=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    no_shadow(c)
    return c

def tbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    return tf

def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=4, space_before=0, line=1.15):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    return p

def run(p, text, size=11, color=TEXTDARK, bold=False, italic=False, font=BODY):
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    return r

# ════════════════════════════════════════════════════════════
def header(slide, title="ANACOLUTHE !?"):
    try:
        slide.shapes.add_picture(IMG + "Logo _anacoluthe_720.jpg",
                                 Inches(0.5), Inches(0.30),
                                 Inches(0.55), Inches(0.55))
    except Exception:
        pass
    tf = tbox(slide, 1.2, 0.30, SW - 2.4, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, title, size=30, color=BRICK, bold=True, font=HEAD)

def footer(slide):
    hline(slide, 0.5, 7.02, SW - 0.5, 7.02)
    tf = tbox(slide, 0.5, 7.06, SW - 1.0, 0.32, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.LEFT)
    run(p, "Anacoluthe !? ", size=8.5, color=BRICK, bold=True, font=HEAD)
    run(p, "· projet BP de Quentin Jaud (promo BSC)", size=8.5, color=TEXTLT, font=HEAD)
    tf2 = tbox(slide, SW - 4.0, 7.06, 3.5, 0.32, anchor=MSO_ANCHOR.MIDDLE)
    p2 = para(tf2, first=True, align=PP_ALIGN.RIGHT)
    run(p2, "V260529 · CC-BY-NC-SA", size=8.5, color=TEXTLT, font=HEAD)

def section_label(slide, x, y, icon, text, color=BLEU):
    # petite etiquette de section (icone + titre uppercase)
    tf = tbox(slide, x, y, SW - 2 * x, 0.36, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    run(p, icon + "  ", size=14, color=color)
    run(p, text, size=12.5, color=color, bold=True, font=HEAD)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — Le but
# ════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1); header(s1); footer(s1)

# Section intro (cadre)
rrect(s1, 0.5, 1.20, SW - 1.0, 1.55, fill=BLOCBG)
section_label(s1, 0.78, 1.05, "\U0001F9ED", "FAIRE ÉQUIPAGE, ÇA S'ANIME AUSSI", color=BLEU)
tf = tbox(s1, 0.8, 1.55, SW - 1.6, 1.15)
p = para(tf, first=True, line=1.4)
run(p, "En embarqué, tu animes déjà la technique, la nav, la sécu. Le ", size=14)
run(p, "« faire équipage »", size=14, bold=True, color=NAVY)
run(p, " (accueillir, répartir les rôles, briefer/débriefer, traverser une tension), c'est ", size=14)
run(p, "l'autre moitié du métier", size=14, bold=True, color=NAVY)
run(p, ".", size=14)

# Revelation (bloc teal prominent)
rrect(s1, 0.5, 3.05, SW - 1.0, 1.55, fill=TEAL50)
rrect(s1, 0.5, 3.05, 0.10, 1.55, fill=TEAL, radius=0.3)  # accent gauche
tf = tbox(s1, 0.85, 3.18, SW - 1.55, 1.3, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, line=1.35, space_after=8)
run(p, "Anacoluthe, c'est un ", size=17, bold=True, color=NAVY)
run(p, "dispositif pédagogique", size=17, bold=True, color=TEAL)
run(p, ", un kit déjà à bord, pour t'épauler là-dessus.", size=17, bold=True, color=NAVY)
p2 = para(tf, line=1.35)
run(p2, "Cartes et affiches sont déjà dans ton bateau, rien à aller chercher.", size=13, color=TEXTDARK)

# Reperes (chips)
reperes = ["déjà dans le bateau", "à la carte", "ça allège ta charge",
           "rien d'obligatoire", "tu restes chef·fe de bord"]
chip_h = 0.42
gap = 0.18
widths = [0.105 * len(t) + 0.45 for t in reperes]
total = sum(widths) + gap * (len(reperes) - 1)
x = (SW - total) / 2.0
cy = 5.05
for t, w in zip(reperes, widths):
    rrect(s1, x, cy, w, chip_h, fill=WHITE, line=NAVY100, line_w=1.0, radius=0.5)
    tf = tbox(s1, x, cy, w, chip_h, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, t, size=10.5, color=TEAL, bold=True, font=HEAD)
    x += w + gap

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Le tableau SOS
# ════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2); header(s2); footer(s2)

# Banniere "prends-le en photo"
rrect(s2, 0.5, 1.12, SW - 1.0, 0.60, fill=AMBER50)
tf = tbox(s2, 0.75, 1.12, SW - 1.5, 0.60, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, line=1.2)
run(p, "\U0001F4F8  Prends-le en photo : ", size=12.5, bold=True, color=AMBER, font=HEAD)
run(p, "ton SOS de poche, à ressortir le moment venu quand ça coince à bord.", size=12, color=TEXTDARK)

# Tableau
rows = [
    ("Quand ça coince à bord…", "Tu pioches…", "Ce que ça débloque"),
    ("L'équipage attend tes consignes pour tout, tu portes l'orga seul·e",
     ("\U0001F465 ", "les rôles tournants", " (bosco, nav, second soigneux, cambusier·ère) + le tableau d'équipage"),
     "chacun·e prend un domaine du jour, l'équipage s'auto-organise, ta charge baisse"),
    ("J1, on reste en surface, attentes non dites",
     ("\U0001F305 ", "le moment accueil + attentes", " (et les accords d'équipage)"),
     "tu sais qui tu as à bord, ton programme colle à leurs envies, des règles co-construites"),
    ("J3-J4, ça frotte, tu deviens l'arbitre",
     ("\U0001F0CF ", "la carte Joker « conflit »", " (ou « rediscuter le cadre »)"),
     "tu nommes la tension sans la trancher seul·e ; le second soigneux du jour facilite"),
    ("Mi-semaine ça patine, un·e décroche",
     ("⚓ ", "le point mi-parcours", ""),
     "on ajuste avant la dernière ligne droite"),
    ("Le débrief du soir ne parle que technique",
     ("\U0001F319 ", "5 min « comment on a bossé ensemble »", " (+ les 5 piliers)"),
     "les acquis (technique ET humain) deviennent nommables, transférables à terre"),
]

tbl_x, tbl_y, tbl_w, tbl_h = 0.5, 1.88, SW - 1.0, 5.05
gframe = s2.shapes.add_table(len(rows), 3, Inches(tbl_x), Inches(tbl_y),
                             Inches(tbl_w), Inches(tbl_h))
table = gframe.table
table.first_row = False
table.horz_banding = False
table.columns[0].width = Inches(4.05)
table.columns[1].width = Inches(4.30)
table.columns[2].width = Inches(tbl_w - 4.05 - 4.30)
table.rows[0].height = Inches(0.42)
for ri in range(1, len(rows)):
    table.rows[ri].height = Inches((tbl_h - 0.42) / (len(rows) - 1))

def cell_setup(cell, fill):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.10)
    cell.margin_top = cell.margin_bottom = Inches(0.04)

# Header row
hdr = rows[0]
for ci, txt in enumerate(hdr):
    cell = table.cell(0, ci)
    cell_setup(cell, BRICK)
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run(p, txt, size=11, color=WHITE, bold=True, font=HEAD)

# Data rows
for ri in range(1, len(rows)):
    left, mid, right = rows[ri]
    fill = WHITE if ri % 2 else BLOCBG
    # col 0
    c0 = table.cell(ri, 0); cell_setup(c0, fill)
    p0 = c0.text_frame.paragraphs[0]; p0.line_spacing = 1.1
    run(p0, left, size=10, color=TEXTDARK)
    # col 1 (emoji + nom en gras + complement)
    c1 = table.cell(ri, 1); cell_setup(c1, fill)
    p1 = c1.text_frame.paragraphs[0]; p1.line_spacing = 1.1
    emoji, strong, rest = mid
    run(p1, emoji, size=11)
    run(p1, strong, size=10, color=NAVY, bold=True)
    if rest:
        run(p1, rest, size=10, color=TEXTDARK)
    # col 2
    c2 = table.cell(ri, 2); cell_setup(c2, fill)
    p2 = c2.text_frame.paragraphs[0]; p2.line_spacing = 1.1
    run(p2, right, size=10, color=TEXTDARK)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — Dispositif complet + invitation
# ════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3); header(s3); footer(s3)

# QR fil-semaine (gros)
qr_size = 2.15
try:
    s3.shapes.add_picture(IMG + "qrcode_anacoluthe_fil_semaine.png",
                          Inches(0.7), Inches(1.35), Inches(qr_size), Inches(qr_size))
except Exception:
    pass
tf = tbox(s3, 0.55, 3.55, 2.45, 0.6)
p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.2)
run(p, "anacoluthe.org/fil-semaine.html", size=9.5, color=BLEU, bold=True, font=HEAD)

# Bloc texte a droite du QR
rrect(s3, 3.25, 1.30, SW - 3.75, 2.30, fill=TEAL50)
section_label(s3, 3.55, 1.15, "\U0001F4F2", "SCANNE LE FIL DE LA SEMAINE", color=TEAL)
tf = tbox(s3, 3.55, 1.65, SW - 4.15, 1.85)
p = para(tf, first=True, line=1.4, space_after=8)
run(p, "Le dispositif complet, déroulé ", size=14, color=TEXTDARK)
run(p, "jour par jour (J1 → J6)", size=14, bold=True, color=NAVY)
run(p, " : quel moment animer quand, quelle carte sortir, quoi dire.", size=14, color=TEXTDARK)
p2 = para(tf, line=1.35)
run(p2, "\U0001F449 Fouille-le ", size=13, color=TEXTDARK)
run(p2, "avant l'arrivée des stagiaires", size=13, bold=True, color=BRICK)
run(p2, ", sur le créneau du samedi, entre le brief et 13h.", size=13, color=TEXTDARK)

# Bande extras (3 colonnes)
ex_y, ex_h = 3.95, 2.95
labels = [
    ("\U0001F392", "Déjà à bord", "Le kit (cartes + affiches plastifiées + guide moniteurice) est déjà dans le bateau."),
    ("\U0001F6E0️", "Ateliers vendredi matin", "« Animer le faire-équipage » et « Fluidifier la semaine avec Anacoluthe »."),
    ("\U0001F4AC", "Communauté WhatsApp", "Pour échanger trucs, retours et coups de main entre moniteurices."),
]
col_w = (SW - 1.0 - 2 * 0.25) / 3.0
cx = 0.5
fills = [BLOCBG, TEAL50, AMBER50]
accents = [NAVY, TEAL, AMBER]
for (icon, title, body), fcol, acc in zip(labels, fills, accents):
    rrect(s3, cx, ex_y, col_w, ex_h, fill=fcol)
    tf = tbox(s3, cx + 0.18, ex_y + 0.18, col_w - 0.36, ex_h - 0.36)
    p = para(tf, first=True, space_after=6)
    run(p, icon + "  ", size=18)
    run(p, title, size=11.5, color=acc, bold=True, font=HEAD)
    p2 = para(tf, line=1.35)
    run(p2, body, size=10.5, color=TEXTDARK)
    cx += col_w + 0.25

# Petit QR WhatsApp dans la 3e colonne
try:
    s3.shapes.add_picture(IMG + "QRcode_anacoluthe_commu-whatsapp.png",
                          Inches(0.5 + 2 * (col_w + 0.25) + col_w - 1.25),
                          Inches(ex_y + ex_h - 1.30),
                          Inches(1.05), Inches(1.05))
except Exception:
    pass

out = "promo/slides_anacoluthe_brief_marseillan.pptx"
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides))
